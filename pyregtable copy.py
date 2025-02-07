import re
import warnings
from collections import Counter
from collections.abc import ValuesView
from typing import Optional, Union

import numpy as np
import pandas as pd
from great_tables import GT
from tabulate import tabulate

from pyfixest.estimation.feiv_ import Feiv
from pyfixest.estimation.feols_ import Feols
from pyfixest.estimation.fepois_ import Fepois
from pyfixest.estimation.FixestMulti_ import FixestMulti
from pyfixest.report.utils import _relabel_expvar
from pyfixest.utils.dev_utils import _select_order_coefs

from docx import Document
from docx.shared import Pt, RGBColor
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

import os

def make_table(
    df: pd.DataFrame,
    type: str = "gt",
    notes: str = "",
    rgroup_sep: str = "tb",
    rgroup_display: bool = True,
    caption: Optional[str] = None,
    tab_label: Optional[str] = None,
    texlocation: str = "htbp",
    full_width: bool = False,
    file_name: Optional[str] = None,
    tab_num: Optional[int] = None,
    **kwargs,
):
    r"""
    Create a booktab style table in the desired format (gt or tex) from a DataFrame.
    The DataFrame can have a multiindex. Column index used to generate horizonal
    table spanners. Row index used to generate row group names and
    row names. The table can have multiple index levels in columns and up to
    two levels in rows.


    Parameters
    ----------
    df : pd.DataFrame
        DataFrame containing the table to be displayed.
    type : str, optional
        Type of table to be created. The default is 'gt'.
    notes : str
        Table notes to be displayed at the bottom of the table.
    rgroup_sep : str
        Whether group names are separated by lines. The default is "tb".
        When output type = 'gt', the options are 'tb', 't', 'b', or '', i.e.
        you can specify whether to have a line above, below, both or none.
        When output type = 'tex' no line will be added between the row groups
        when rgroup_sep is '' and otherwise a line before the group name will be added.
    rgroup_display : bool
        Whether to display row group names. The default is
        True.
    caption : str
        Table caption to be displayed at the top of the table. The default is None.
        When either caption or label is provided the table will be wrapped in a
        table environment.
    tab_label : str
        LaTex label of the table. The default is None. When either caption or label
        is provided the table will be wrapped in a table environment.
    texlocation : str
        Location of the table. The default is 'htbp'.
    full_width : bool
        Whether to expand the table to the full width of the page. The default is False.
    file_name : str
        Name of the file to save the table to. The default is None.
        gt tables will be saved as html files and latex tables as tex files.
    tab_num : int
        Table number to be replaced in the docx document. The default is None.
        When a  positive number is provided the table at position tab_num will be 
        replaced by the new table in the document. If tab_num is larger than 
        the number of tables in the document the new table will be appended to the end.
        Only supported for docx output.

    Returns
    -------
    A table in the specified format.
    """
    assert isinstance(df, pd.DataFrame), "df must be a pandas DataFrame."
    assert not isinstance(df.index, pd.MultiIndex) or df.index.nlevels <= 2, (
        "Row index can have at most two levels."
    )
    assert type in ["gt", "tex", "docx", "pptx"], "type must be either 'gt', 'tex' or 'docx'."
    assert rgroup_sep in [
        "tb",
        "t",
        "b",
        "",
    ], "rgroup_sep must be either 'tb', 't', 'b', or ''."
    assert file_name is None or (
        isinstance(file_name, str) and file_name.endswith((".html", ".tex", ".docx", "pptx"))
    ), "file_name must end with '.html', '.tex', or '.docx'."
    # check: tab_num currently only supported for docx output
    assert tab_num is None or type == "docx" or type == "pptx", "tab_num currently only supported for docx output"
    # Make a copy of the DataFrame to avoid modifying the original
    dfs = df.copy()

    # Produce LaTeX code if either type is 'tex' or the
    # user has passed a file_name which ends with '.tex'
    if type == "tex" or (isinstance(file_name, str) and file_name.endswith(".tex")):
        # First wrap all cells which contain a line break in a makecell command
        dfs = dfs.map(
            lambda x: f"\\makecell{{{x}}}" if isinstance(x, str) and "\\\\" in x else x
        )
        row_levels = dfs.index.nlevels
        # when the row index has more than one level, we will store
        # the top level to use later to add clines and row group titles
        # and then remove it
        if row_levels > 1:
            # Store the top level of the row index
            top_row_id = dfs.index.get_level_values(0).to_list()
            # Generate a list of the distinct values
            row_groups = list(dict.fromkeys(top_row_id))
            # Generate a list containing the number of rows for each group
            row_groups_len = [top_row_id.count(group) for group in row_groups]
            # Drop the top level of the row index:
            dfs.index = dfs.index.droplevel(0)

        # Style the table
        styler = dfs.style
        # if caption is not None:
        #     styler.set_caption(caption)

        # Generate LaTeX code
        latex_res = styler.to_latex(
            hrules=True,
            multicol_align="c",
            multirow_align="t",
            column_format="l" + "c" * (dfs.shape[1] + dfs.index.nlevels),
        )

        # # Now perform post-processing of the LaTeX code
        # # First split the LaTeX code into lines
        lines = latex_res.splitlines()
        # Find the line number of the \midrule
        line_at = next(i for i, line in enumerate(lines) if "\\midrule" in line)
        # Add space after this \midrule:
        lines.insert(line_at + 1, "\\addlinespace")
        line_at += 1

        # When there are row groups then insert midrules and groupname
        if row_levels > 1 and len(row_groups) > 1:
            # Insert a midrule after each row group
            for i in range(len(row_groups)):
                if rgroup_display:
                    # Insert a line with the row group name & same space around it
                    # lines.insert(line_at+1, "\\addlinespace")
                    lines.insert(line_at + 1, "\\emph{" + row_groups[i] + "} \\\\")
                    lines.insert(line_at + 2, "\\addlinespace")
                    lines.insert(line_at + 3 + row_groups_len[i], "\\addlinespace")
                    line_at += 3
                if (rgroup_sep != "") and (i < len(row_groups) - 1):
                    # For tex output we only either at a line between the row groups or not
                    # And we don't add a line after the last row group
                    line_at += row_groups_len[i] + 1
                    lines.insert(line_at, "\\midrule")
                    lines.insert(line_at + 1, "\\addlinespace")
                    line_at += 1
        else:
            # Add line space before the end of the table
            lines.insert(line_at + dfs.shape[0] + 1, "\\addlinespace")

        # Insert cmidrules (equivalent to column spanners in gt)
        # First find the first line with an occurrence of "multicolumn"
        cmidrule_line_number = None
        for i, line in enumerate(lines):
            if "multicolumn" in line:
                cmidrule_line_number = i + 1
                # Regular expression to find \multicolumn{number}
                pattern = r"\\multicolumn\{(\d+)\}"
                # Find all matches (i.e. values of d) in the LaTeX string & convert to integers
                ncols = [int(match) for match in re.findall(pattern, line)]

                cmidrule_string = ""
                leftcol = 2
                for n in ncols:
                    cmidrule_string += (
                        r"\cmidrule(lr){"
                        + str(leftcol)
                        + "-"
                        + str(leftcol + n - 1)
                        + "} "
                    )
                    leftcol += n
                lines.insert(cmidrule_line_number, cmidrule_string)

        # # Put the lines back together
        latex_res = "\n".join(lines)

        # Wrap in threeparttable to allow for table notes
        if notes is not None:
            latex_res = (
                "\\begin{threeparttable}\n"
                + latex_res
                + "\n\\footnotesize "
                + notes
                + "\n\\end{threeparttable}"
            )
        else:
            latex_res = (
                "\\begin{threeparttable}\n" + latex_res + "\n\\end{threeparttable}"
            )

        # If caption or label specified then wrap in table environment
        if (caption is not None) or (tab_label is not None):
            latex_res = (
                "\\begin{table}["
                + texlocation
                + "]\n"
                + "\\centering\n"
                + ("\\caption{" + caption + "}\n" if caption is not None else "")
                + ("\\label{" + tab_label + "}\n" if tab_label is not None else "")
                + latex_res
                + "\n\\end{table}"
            )

        # Set cell aligment to top
        latex_res = "\\renewcommand\\cellalign{t}\n" + latex_res

        # Set table width to full page width if full_width is True
        # This is done by changing the tabular environment to tabular*
        if full_width:
            latex_res = latex_res.replace(
                "\\begin{tabular}{l", "\\begin{tabularx}{\\linewidth}{X"
            )
            latex_res = latex_res.replace(
                "\\end{tabular}", "\\end{tabularx}\n \\vspace{3pt}"
            )
            # with tabular*
            # latex_res = latex_res.replace("\\begin{tabular}{", "\\begin{tabular*}{\linewidth}{@{\extracolsep{\\fill}}")
            # latex_res = latex_res.replace("\\end{tabular}", "\\end{tabular*}")

        if file_name is not None:
            with open(file_name, "w") as f:
                f.write(latex_res)  # Write the latex code to a file

        if type == "tex":
            return latex_res

    if type == "gt":
        # GT does not support MultiIndex columns, so we need to flatten the columns
        if isinstance(dfs.columns, pd.MultiIndex):
            # Store labels of the last level of the column index (to use as column names)
            col_names = dfs.columns.get_level_values(-1)
            nl = dfs.columns.nlevels
            # As GT does not accept non-unique column names: so to allow for them
            # we just assign column numbers to the lowest index level
            col_numbers = list(map(str, range(len(dfs.columns))))
            # Save the whole column index in order to generate table spanner labels later
            dfcols = dfs.columns.to_list()
            # Then flatten the column index just numbering the columns
            dfs.columns = pd.Index(col_numbers)
            # Store the mapping of column numbers to column names
            col_dict = dict(zip(col_numbers, col_names))
            # Modify the last elements in each tuple in dfcols
            dfcols = [(t[:-1] + (col_numbers[i],)) for i, t in enumerate(dfcols)]
            # And drop the first column as we don't want table spanners on top of the variables
            # WE DON'T NEED THIS WITH ROW INDEX dfcols = dfcols[1:]
        else:
            nl = 1

        rowindex = dfs.index

        # Now reset row index to have the index as columns to be displayed in the table
        dfs.reset_index(inplace=True)

        # And specify the rowname_col and groupname_col
        if isinstance(rowindex, pd.MultiIndex):
            rowname_col = dfs.columns[1]
            groupname_col = dfs.columns[0]
        else:
            rowname_col = dfs.columns[0]
            groupname_col = None

        # Generate the table with GT
        gt = GT(dfs, auto_align=False)

        # When caption is provided, add it to the table
        if caption is not None:
            gt = (
                gt.tab_header(title=caption).tab_options(
                    table_border_top_style="hidden",
                )  # Otherwise line above caption
            )

        if nl > 1:
            # Add column spanners based on multiindex
            # Do this for every level in the multiindex (except the one with the column numbers)
            for i in range(nl - 1):
                col_spanners: dict[str, list[str | int]] = {}
                # Iterate over columns and group them by the labels in the respective level
                for c in dfcols:
                    key = c[i]
                    if key not in col_spanners:
                        col_spanners[key] = []
                    col_spanners[key].append(c[-1])
                for label, columns in col_spanners.items():
                    gt = gt.tab_spanner(label=label, columns=columns, level=nl - 1 - i)
            # Restore column names
            gt = gt.cols_label(**col_dict)

        # Customize the table layout
        gt = (
            gt.tab_source_note(notes)
            .tab_stub(rowname_col=rowname_col, groupname_col=groupname_col)
            .tab_options(
                table_border_bottom_style="hidden",
                stub_border_style="hidden",
                column_labels_border_top_style="solid",
                column_labels_border_top_color="black",
                column_labels_border_bottom_style="solid",
                column_labels_border_bottom_color="black",
                column_labels_border_bottom_width="0.5px",
                column_labels_vlines_color="white",
                column_labels_vlines_width="0px",
                table_body_border_top_style="solid",
                table_body_border_top_width="0.5px",
                table_body_border_top_color="black",
                table_body_hlines_style="none",
                table_body_vlines_color="white",
                table_body_vlines_width="0px",
                table_body_border_bottom_color="black",
                row_group_border_top_style="solid",
                row_group_border_top_width="0.5px",
                row_group_border_top_color="black",
                row_group_border_bottom_style="solid",
                row_group_border_bottom_width="0.5px",
                row_group_border_bottom_color="black",
                row_group_border_left_color="white",
                row_group_border_right_color="white",
                data_row_padding="4px",
                column_labels_padding="4px",
            )
            .cols_align(align="center")
        )

        # Full page width
        if full_width:
            gt = gt.tab_options(table_width="100%")

        # Customize row group display
        if "t" not in rgroup_sep:
            gt = gt.tab_options(row_group_border_top_style="none")
        if "b" not in rgroup_sep:
            gt = gt.tab_options(row_group_border_bottom_style="none")
        if not rgroup_display:
            gt = gt.tab_options(
                row_group_font_size="0px",
                row_group_padding="0px",
            )
        # Save the html code of the table to a file
        if file_name is not None:
            with open(file_name, "w") as f:
                f.write(gt.as_raw_html())

        return gt

    # Produce MS Word table if type is 'docx' or the user has passed a file_name which ends with '.docx'
    if type == "docx" or (isinstance(file_name, str) and file_name.endswith(".docx")):

        # Number of headline levels
        headline_levels = dfs.columns.nlevels
        # Are there row groups: is the case when dfs.index.nlevels > 1
        row_groups = (dfs.index.nlevels>1)
        # Number of columns
        ncols = dfs.shape[1] +1

        # Check if the document exists
        if file_name and os.path.exists(file_name):
            document = Document(file_name)
            # Determine the number of tables in the document
            n_tables = len(document.tables)
        else:
            document = Document()

        if n_tables is not None and n_tables > 0 and tab_num is not None and tab_num <= n_tables:
            # Replace the table at position tab_num
            table = document.tables[tab_num - 1]
            # Replace the caption before the table
            if caption is not None:
                # Find the paragraph before the table
                table_idx = list(document._body._body).index(table._element)
                if table_idx > 0:
                    prev_par_element = document._body._element[table_idx - 1]
                    if prev_par_element.tag.endswith('p') and 'Table' in prev_par_element.text:
                        # replace text in last subelement of prev_par_element (this should be the old caption)
                        prev_par_element[-1].text = f': {caption}'
            # Delete all rows in the old table
            for row in table.rows:
                table._element.remove(row._element)
        else:
            # Add caption and new table
            if caption is not None:
                paragraph = document.add_paragraph('Table ', style='Caption')
                run = paragraph.add_run()
                r = run._r
                fldChar = OxmlElement('w:fldChar')
                fldChar.set(qn('w:fldCharType'), 'begin')
                r.append(fldChar)
                instrText = OxmlElement('w:instrText')
                instrText.text = r'SEQ Table \* ARABIC'
                r.append(instrText)
                fldChar = OxmlElement('w:fldChar')
                fldChar.set(qn('w:fldCharType'), 'end')
                r.append(fldChar)
                bold_run = paragraph.add_run(f': {caption}')
                bold_run.bold = False
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                # Set the font color to black and font size to 11
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    run.font.size = Pt(11)

            table = document.add_table(rows=0, cols=ncols)
            table.style = 'Table Grid'

        # Add column headers
        if isinstance(dfs.columns, pd.MultiIndex):
            # Add multiple headline rows for MultiIndex columns
            for level in range(headline_levels):
                hdr_cells = table.add_row().cells
                prev_col = None
                prev_cell_index = None
                for i, col in enumerate(dfs.columns.get_level_values(level)):
                    cell_index = i + 1
                    if col != prev_col:
                        hdr_cells[cell_index].text = str(col)
                        prev_col = col
                        prev_cell_index = cell_index
                    else:
                        hdr_cells[prev_cell_index].merge(hdr_cells[cell_index])
        else:
            hdr_cells = table.add_row().cells
            for i, col in enumerate(dfs.columns):
                hdr_cells[i + 1].text = str(col)

        # Add row names and data
        row_group_rows=[]
        if row_groups:
            current_group = None
            for idx, row in dfs.iterrows():
                if idx[0] != current_group:
                    # New row group
                    current_group = idx[0]
                    # append row number to row_group_rows
                    row_group_rows.append(len(table.rows))
                    if rgroup_display:
                        # Add a row for the group name
                        group_row_cells = table.add_row().cells
                        # add row group name
                        group_row_cells[0].text = str(current_group)
                        # make this cell slightly taller
                        for paragraph in group_row_cells[0].paragraphs:
                            paragraph.paragraph_format.space_after = Pt(3)
                            paragraph.paragraph_format.space_before = Pt(3)
                        for cell in group_row_cells[1:]:
                            cell.text = ""
                row_cells = table.add_row().cells
                row_cells[0].text = str(idx[1])
                for i, val in enumerate(row):
                    row_cells[i + 1].text = str(val)
        else:
            for idx, row in dfs.iterrows():
                row_cells = table.add_row().cells
                row_cells[0].text = str(idx)
                for i, val in enumerate(row):
                    row_cells[i + 1].text = str(val)

        # Center all columns except the first one
        for row in table.rows:
            for cell in row.cells[1:]:
                for paragraph in cell.paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add notes if provided
        if notes:
            #document.add_paragraph(notes)
            # Add row to the table that consists only of one cell with the notes
            notes_row = table.add_row().cells
            notes_row[0].text = notes
            # Merge the cell with the notes
            table.cell(-1, 0).merge(table.cell(-1, ncols-1))
            # Set alignment and font size for the notes
            for paragraph in notes_row[0].paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                for run in paragraph.runs:
                    run.font.size = Pt(9)
 
        # First hide all borders 
        for row in table.rows:
            for cell in row.cells:
                tcPr = cell._element.get_or_add_tcPr()
                borders = OxmlElement('w:tcBorders')
                top_border = OxmlElement('w:top')
                top_border.set(qn('w:val'), 'nil')
                borders.append(top_border)
                bottom_border = OxmlElement('w:bottom')
                bottom_border.set(qn('w:val'), 'nil')
                borders.append(bottom_border)
                left_border = OxmlElement('w:left')
                left_border.set(qn('w:val'), 'nil')
                borders.append(left_border)
                right_border = OxmlElement('w:right')
                right_border.set(qn('w:val'), 'nil')
                borders.append(right_border)
                tcPr.append(borders)

        # Add a thicker line above the top row
        for cell in table.rows[0].cells:
            tcPr = cell._element.get_or_add_tcPr()
            borders = OxmlElement('w:tcBorders')
            top_border = OxmlElement('w:top')
            top_border.set(qn('w:val'), 'single')
            top_border.set(qn('w:sz'), '8')
            borders.append(top_border)
            tcPr.append(borders)

        # Add a line to the last headline row
        for cell in table.rows[headline_levels-1].cells:
            tcPr = cell._element.get_or_add_tcPr()
            borders = OxmlElement('w:tcBorders')
            bottom_border = OxmlElement('w:bottom')
            bottom_border.set(qn('w:val'), 'single')
            bottom_border.set(qn('w:sz'), '4')
            borders.append(bottom_border)
            tcPr.append(borders)

        # If the column index has more than one level, add a line above the last headline 
        # row that spanns all but the first cell
        if headline_levels > 1:
            for cell in table.rows[headline_levels-1].cells[1:]:
                tcPr = cell._element.get_or_add_tcPr()
                borders = OxmlElement('w:tcBorders')
                top_border = OxmlElement('w:top')
                top_border.set(qn('w:val'), 'single')
                top_border.set(qn('w:sz'), '4')
                borders.append(top_border)
                tcPr.append(borders)

        # Loop over  all lines in row_group_rows
        # And add lines above and below the row group names depending on rgroup_display and rgroup_sep
        for row in row_group_rows:
            if "t" in rgroup_sep:
                # Add a line above the row group name
                for cell in table.rows[row].cells:
                    tcPr = cell._element.get_or_add_tcPr()
                    borders = OxmlElement('w:tcBorders')
                    top_border = OxmlElement('w:top')
                    top_border.set(qn('w:val'), 'single')
                    top_border.set(qn('w:sz'), '4')
                    borders.append(top_border)
                    tcPr.append(borders)
            if rgroup_display and "b" in rgroup_sep:
                # Add a line below the row group name
                for cell in table.rows[row].cells:
                    tcPr = cell._element.get_or_add_tcPr()
                    borders = OxmlElement('w:tcBorders')
                    bottom_border = OxmlElement('w:bottom')
                    bottom_border.set(qn('w:val'), 'single')
                    bottom_border.set(qn('w:sz'), '4')
                    borders.append(bottom_border)
                    tcPr.append(borders)

        # Add a thicker line below the last row
        for cell in table.rows[-2].cells:
            tcPr = cell._element.get_or_add_tcPr()
            borders = OxmlElement('w:tcBorders')
            bottom_border = OxmlElement('w:bottom')
            bottom_border.set(qn('w:val'), 'single')
            bottom_border.set(qn('w:sz'), '8')
            borders.append(bottom_border)
            tcPr.append(borders)

        # Adapt cell margins
        tc = table._element
        tblPr = tc.tblPr
        tblCellMar = OxmlElement('w:tblCellMar')
        # set left and right margins to zero
        # and top and bottom margins to 60 dxa
        kwargs = {"left":0, "right":0, "top":60, "bottom":60}
        for m in ["left","right","top","bottom"]:
            node = OxmlElement("w:{}".format(m))
            node.set(qn('w:w'), str(kwargs.get(m)))
            node.set(qn('w:type'), 'dxa')
            tblCellMar.append(node)
        tblPr.append(tblCellMar)

        # Save the document
        if file_name is not None:
            document.save(file_name)
        if type == "docx":
            return document

    if type == "pptx" or (isinstance(file_name, str) and file_name.endswith(".pptx")):

        # Number of headline levels
        headline_levels = dfs.columns.nlevels
        # Are there row groups: is the case when dfs.index.nlevels > 1
        row_groups = (dfs.index.nlevels>1)
        # Number of columns
        ncols = dfs.shape[1] +1

        # Create a new presentation or open an existing one
        if file_name and os.path.exists(file_name):
            presentation = Presentation(file_name)
        else:
            presentation = Presentation()

        # Add a slide with a title and content layout
        slide_layout = presentation.slide_layouts[5]  # Use a blank layout
        slide = presentation.slides.add_slide(slide_layout)

        # Add title to the slide
        if caption is not None:
            title = slide.shapes.title
            title.text = caption

        # Add a table to the slide
        rows, cols = dfs.shape
        table = slide.shapes.add_table(rows + 1, cols + 1, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table

        # Set column headers
        for i, col in enumerate(dfs.columns):
            table.cell(0, i + 1).text = str(col)
            # Add column headers
            if isinstance(dfs.columns, pd.MultiIndex):
                # Add multiple headline rows for MultiIndex columns
                for level in range(headline_levels):
                    hdr_cells = table.add_row().cells
                    prev_col = None
                    prev_cell_index = None
                    for i, col in enumerate(dfs.columns.get_level_values(level)):
                        cell_index = i + 1
                        if col != prev_col:
                            hdr_cells[cell_index].text = str(col)
                            prev_col = col
                            prev_cell_index = cell_index
                        else:
                            hdr_cells[prev_cell_index].merge(hdr_cells[cell_index])
            else:
                hdr_cells = table.add_row().cells
                for i, col in enumerate(dfs.columns):
                    hdr_cells[i + 1].text = str(col)

        # Set row headers and data
        for i, idx in enumerate(dfs.index):
            table.cell(i + 1, 0).text = str(idx)
            for j, val in enumerate(dfs.iloc[i]):
                table.cell(i + 1, j + 1).text = str(val)

        # Format the table
        for row in table.rows:
            for cell in row.cells:
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Save the presentation
        if file_name is not None:
            presentation.save(file_name)
        if type == "pptx":
            return presentation

def _relabel_index(index, labels=None, stats_labels=None):
    if stats_labels is None:
        if isinstance(index, pd.MultiIndex):
            index = pd.MultiIndex.from_tuples(
                [tuple(labels.get(k, k) for k in i) for i in index]
            )
        else:
            index = [labels.get(k, k) for k in index]
    else:
        # if stats_labels is provided, we relabel the lowest level of the index with it
        if isinstance(index, pd.MultiIndex):
            new_index = []
            for i in index:
                new_index.append(
                    tuple(
                        [labels.get(k, k) for k in i[:-1]]
                        + [stats_labels.get(i[-1], i[-1])]
                    )
                )
            index = pd.MultiIndex.from_tuples(new_index)
        else:
            index = [stats_labels.get(k, k) for k in index]
    return index

