import re
import warnings
from collections import Counter
from collections.abc import ValuesView
from typing import Optional, Union

import numpy as np
import pandas as pd
from great_tables import GT
from tabulate import tabulate

from pyfixest.estimation.feiv_ import Feiv
from pyfixest.estimation.feols_ import Feols
from pyfixest.estimation.fepois_ import Fepois
from pyfixest.estimation.FixestMulti_ import FixestMulti
from pyfixest.report.utils import _relabel_expvar
from pyfixest.utils.dev_utils import _select_order_coefs

from docx import Document
from docx.shared import Pt, RGBColor
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

import os

def _relabel_index(index, labels=None, stats_labels=None):
    if stats_labels is None:
        if isinstance(index, pd.MultiIndex):
            index = pd.MultiIndex.from_tuples(
                [tuple(labels.get(k, k) for k in i) for i in index]
            )
        else:
            index = [labels.get(k, k) for k in index]
    else:
        # if stats_labels is provided, we relabel the lowest level of the index with it
        if isinstance(index, pd.MultiIndex):
            new_index = []
            for i in index:
                new_index.append(
                    tuple(
                        [labels.get(k, k) for k in i[:-1]]
                        + [stats_labels.get(i[-1], i[-1])]
                    )
                )
            index = pd.MultiIndex.from_tuples(new_index)
        else:
            index = [stats_labels.get(k, k) for k in index]
    return index



class TabOut:
    def __init__(self, df: pd.DataFrame, 
                 notes: str = "", 
                 caption: Optional[str] = None, 
                 tab_label: Optional[str] = None,
                 rgroup_sep: str = "tb",
                 rgroup_display: bool = True,
                ):
        assert isinstance(df, pd.DataFrame), "df must be a pandas DataFrame."
        assert not isinstance(df.index, pd.MultiIndex) or df.index.nlevels <= 2, (
            "Row index can have at most two levels."
        )
        self.df = df
        self.notes = notes
        self.caption = caption
        self.tab_label = tab_label
        self.rgroup_sep = rgroup_sep
        self.rgroup_display = rgroup_display
    
    def _revise_file_path(self, file_path: str, label: str, extension: str) -> str:
        if extension=='gt':
            # gt is always saved as html
            extension = 'html'
        if file_path.endswith(extension):
            # if the path ends with the specified extension, we return it
            return file_path
        elif '.' in file_path:
            # if it has another file extension, we replace it with the specified extension 
            return file_path.rsplit('.', 1)[0] + extension
        elif os.path.isdir(file_path):
            # otherwise we check whether the file_path is a directory, 
            # and define a new file name in this directory named label.extension
            return os.path.join(file_path, f"{label}.{extension}")
        else:
            # otherwise we add the extension to the file_path
            return f"{file_path}.{extension}"
            
    def make(self, types: Union[str, list[str]], path: Optional[str] = None, **kwargs):
        if isinstance(types, str):
            types = [types]
        assert all(t in ["gt","tex", "docx", "pptx"] for t in types), "types must be either 'gt', 'tex', 'docx', or 'pptx'."
        assert path is None or (
            isinstance(path, str) and (path.endswith((".html", ".tex", ".docx", ".pptx")) or '.' not in path)
        ), "path must end with '.html', '.tex', '.docx', '.pptx', or have no file extension."
        # MOVE TO WORD OUTPUT check: tab_num currently only supported for docx and pptx output
        # assert tab_num is None or any(t in ["docx", "pptx"] for t in types), "tab_num currently only supported for docx and pptx output"

        results = {}
        for t in types:
            if path is not None:
                file_name = self._revise_file_path(path, self.tab_label, t)
            else: 
                file_name = None
            if t == "gt":
                results[t] = self._output_gt(file_name=file_name, **kwargs)
            if t == "tex":
                results[t] = self._output_tex(file_name=file_name, **kwargs)
            if t == "docx":
                results[t] = self._output_docx(file_name=file_name, **kwargs)
            if t == "pptx":
                results[t] = self._output_pptx(file_name=file_name, **kwargs)
        # Return only the first output type result in the list
        return results[types[0]]
    
    def _output_pptx(self, file_name: Optional[str] = None, **kwargs):
        # Make a copy of the DataFrame to avoid modifying the original
        dfs = self.df.copy()

        # Number of headline levels
        headline_levels = dfs.columns.nlevels
        # Are there row groups: is the case when dfs.index.nlevels > 1
        row_groups = (dfs.index.nlevels > 1)
        if row_groups:
            # Determine number of row groups
            nrow_groups = len(dfs.index.get_level_values(0).unique())
        # Number of rows (excluding headline rows)
        nrows = dfs.shape[0]
        # Number of columns
        ncols = dfs.shape[1] + 1

        # Create a new presentation or open an existing one
        if file_name and os.path.exists(file_name):
            presentation = Presentation(file_name)
        else:
            presentation = Presentation()

        # Add a slide with a title and content layout
        slide_layout = presentation.slide_layouts[5]  # Use a blank layout
        slide = presentation.slides.add_slide(slide_layout)

        # Add title to the slide
        if self.caption is not None:
            title = slide.shapes.title
            title.text = self.caption

        # Add a table to the slide
        if self.rgroup_display:
            table = slide.shapes.add_table(nrows + headline_levels+ nrow_groups, ncols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
        else:
            table = slide.shapes.add_table(nrows + headline_levels, ncols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table
        # Set column headers
        if isinstance(dfs.columns, pd.MultiIndex):
            # Add multiple headline rows for MultiIndex columns
            for level in range(headline_levels):
                prev_val = None
                start_cell = None
                for i, val in enumerate(dfs.columns.get_level_values(level)):
                    if val != prev_val:
                        if start_cell is not None:
                            # merge prior cells
                            start_cell.merge(table.cell(level, i))
                        # assign index value to cell
                        table.cell(level, i+1).text = str(val)
                        # and define this as a new start_cell
                        start_cell= table.cell(level, i+1)
                    else:
                        # if the value is the same as the previous one, we just leave the cell empty
                        table.cell(level, i+1).text = ""
                        # if it is the last cell in the row, we merge it with the previous cell
                        if i == len(dfs.columns.get_level_values(level)) - 1:
                            start_cell.merge(table.cell(level, i+1))
                    prev_val = val
        else:
            for i, col in enumerate(dfs.columns):
                table.cell(0, i + 1).text = str(col)

        # Add row names and data
        if row_groups:
            current_group = None
            row_num=headline_levels
            for i, idx in enumerate(dfs.index):
                if idx[0] != current_group:
                    # New row group
                    current_group = idx[0]
                    if self.rgroup_display:
                        # Add group name
                        table.cell(row_num, 0).text = str(current_group)
                        row_num += 1
                table.cell(row_num, 0).text = str(idx[1])
                for j, val in enumerate(dfs.iloc[i]):
                    table.cell(row_num, j + 1).text = str(val)
                row_num += 1
        else:
            for i, idx in enumerate(dfs.index):
                table.cell(i + headline_levels, 0).text = str(idx)
                for j, val in enumerate(dfs.iloc[i]):
                    table.cell(i + headline_levels, j + 1).text = str(val)

        # Format the table
        for row in table.rows:
            for cell in row.cells:
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                #cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Save the presentation
        if file_name is not None:
            presentation.save(file_name)

        return presentation

    def _output_docx(self, file_name: Optional[str] = None, tab_num: Optional[int] = None, **kwargs):
        # Make a copy of the DataFrame to avoid modifying the original
        dfs = self.df.copy()

        # Number of headline levels
        headline_levels = dfs.columns.nlevels
        # Are there row groups: is the case when dfs.index.nlevels > 1
        row_groups = (dfs.index.nlevels > 1)
        # Number of columns
        ncols = dfs.shape[1] + 1

        # Check if the document exists
        if file_name and os.path.exists(file_name):
            document = Document(file_name)
            # Determine the number of tables in the document
            n_tables = len(document.tables)
        else:
            document = Document()
            n_tables = 0

        if n_tables > 0 and tab_num is not None and tab_num <= n_tables:
            # Replace the table at position tab_num
            table = document.tables[tab_num - 1]
            # Replace the caption before the table
            if self.caption is not None:
                # Find the paragraph before the table
                table_idx = list(document._body._body).index(table._element)
                if table_idx > 0:
                    prev_par_element = document._body._element[table_idx - 1]
                    if prev_par_element.tag.endswith('p') and 'Table' in prev_par_element.text:
                        # replace text in last subelement of prev_par_element (this should be the old caption)
                        prev_par_element[-1].text = f': {self.caption}'
            # Delete all rows in the old table
            for row in table.rows:
                table._element.remove(row._element)
        else:
            # Add caption and new table
            if self.caption is not None:
                paragraph = document.add_paragraph('Table ', style='Caption')
                run = paragraph.add_run()
                r = run._r
                fldChar = OxmlElement('w:fldChar')
                fldChar.set(qn('w:fldCharType'), 'begin')
                r.append(fldChar)
                instrText = OxmlElement('w:instrText')
                instrText.text = r'SEQ Table \* ARABIC'
                r.append(instrText)
                fldChar = OxmlElement('w:fldChar')
                fldChar.set(qn('w:fldCharType'), 'end')
                r.append(fldChar)
                bold_run = paragraph.add_run(f': {self.caption}')
                bold_run.bold = False
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                # Set the font color to black and font size to 11
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    run.font.size = Pt(11)

            table = document.add_table(rows=0, cols=ncols)
            table.style = 'Table Grid'

        # Add column headers
        if isinstance(dfs.columns, pd.MultiIndex):
            # Add multiple headline rows for MultiIndex columns
            for level in range(headline_levels):
                hdr_cells = table.add_row().cells
                prev_col = None
                prev_cell_index = None
                for i, col in enumerate(dfs.columns.get_level_values(level)):
                    cell_index = i + 1
                    if col != prev_col:
                        hdr_cells[cell_index].text = str(col)
                        prev_col = col
                        prev_cell_index = cell_index
                    else:
                        hdr_cells[prev_cell_index].merge(hdr_cells[cell_index])
        else:
            hdr_cells = table.add_row().cells
            for i, col in enumerate(dfs.columns):
                hdr_cells[i + 1].text = str(col)

        # Add row names and data
        row_group_rows = []
        if row_groups:
            current_group = None
            for idx, row in dfs.iterrows():
                if idx[0] != current_group:
                    # New row group
                    current_group = idx[0]
                    # append row number to row_group_rows
                    row_group_rows.append(len(table.rows))
                    if self.rgroup_display:
                        # Add a row for the group name
                        group_row_cells = table.add_row().cells
                        # add row group name
                        group_row_cells[0].text = str(current_group)
                        # make this cell slightly taller
                        for paragraph in group_row_cells[0].paragraphs:
                            paragraph.paragraph_format.space_after = Pt(3)
                            paragraph.paragraph_format.space_before = Pt(3)
                        for cell in group_row_cells[1:]:
                            cell.text = ""
                row_cells = table.add_row().cells
                row_cells[0].text = str(idx[1])
                for i, val in enumerate(row):
                    row_cells[i + 1].text = str(val)
        else:
            for idx, row in dfs.iterrows():
                row_cells = table.add_row().cells
                row_cells[0].text = str(idx)
                for i, val in enumerate(row):
                    row_cells[i + 1].text = str(val)

        # Center all columns except the first one
        for row in table.rows:
            for cell in row.cells[1:]:
                for paragraph in cell.paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add notes if provided
        if self.notes:
            # Add row to the table that consists only of one cell with the notes
            notes_row = table.add_row().cells
            notes_row[0].text = self.notes
            # Merge the cell with the notes
            table.cell(-1, 0).merge(table.cell(-1, ncols - 1))
            # Set alignment and font size for the notes
            for paragraph in notes_row[0].paragraphs:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                for run in paragraph.runs:
                    run.font.size = Pt(9)

        # First hide all borders
        for row in table.rows:
            for cell in row.cells:
                tcPr = cell._element.get_or_add_tcPr()
                borders = OxmlElement('w:tcBorders')
                top_border = OxmlElement('w:top')
                top_border.set(qn('w:val'), 'nil')
                borders.append(top_border)
                bottom_border = OxmlElement('w:bottom')
                bottom_border.set(qn('w:val'), 'nil')
                borders.append(bottom_border)
                left_border = OxmlElement('w:left')
                left_border.set(qn('w:val'), 'nil')
                borders.append(left_border)
                right_border = OxmlElement('w:right')
                right_border.set(qn('w:val'), 'nil')
                borders.append(right_border)
                tcPr.append(borders)

        # Add a thicker line above the top row
        for cell in table.rows[0].cells:
            tcPr = cell._element.get_or_add_tcPr()
            borders = OxmlElement('w:tcBorders')
            top_border = OxmlElement('w:top')
            top_border.set(qn('w:val'), 'single')
            top_border.set(qn('w:sz'), '8')
            borders.append(top_border)
            tcPr.append(borders)

        # Add a line to the last headline row
        for cell in table.rows[headline_levels - 1].cells:
            tcPr = cell._element.get_or_add_tcPr()
            borders = OxmlElement('w:tcBorders')
            bottom_border = OxmlElement('w:bottom')
            bottom_border.set(qn('w:val'), 'single')
            bottom_border.set(qn('w:sz'), '4')
            borders.append(bottom_border)
            tcPr.append(borders)

        # If the column index has more than one level, add a line above the last headline
        # row that spans all but the first cell
        if headline_levels > 1:
            for cell in table.rows[headline_levels - 1].cells[1:]:
                tcPr = cell._element.get_or_add_tcPr()
                borders = OxmlElement('w:tcBorders')
                top_border = OxmlElement('w:top')
                top_border.set(qn('w:val'), 'single')
                top_border.set(qn('w:sz'), '4')
                borders.append(top_border)
                tcPr.append(borders)

        # Loop over all lines in row_group_rows
        # And add lines above and below the row group names depending on rgroup_display and rgroup_sep
        for row in row_group_rows:
            if "t" in self.rgroup_sep:
                # Add a line above the row group name
                for cell in table.rows[row].cells:
                    tcPr = cell._element.get_or_add_tcPr()
                    borders = OxmlElement('w:tcBorders')
                    top_border = OxmlElement('w:top')
                    top_border.set(qn('w:val'), 'single')
                    top_border.set(qn('w:sz'), '4')
                    borders.append(top_border)
                    tcPr.append(borders)
            if self.rgroup_display and "b" in self.rgroup_sep:
                # Add a line below the row group name
                for cell in table.rows[row].cells:
                    tcPr = cell._element.get_or_add_tcPr()
                    borders = OxmlElement('w:tcBorders')
                    bottom_border = OxmlElement('w:bottom')
                    bottom_border.set(qn('w:val'), 'single')
                    bottom_border.set(qn('w:sz'), '4')
                    borders.append(bottom_border)
                    tcPr.append(borders)

        # Add a thicker line below the last row
        for cell in table.rows[-2].cells:
            tcPr = cell._element.get_or_add_tcPr()
            borders = OxmlElement('w:tcBorders')
            bottom_border = OxmlElement('w:bottom')
            bottom_border.set(qn('w:val'), 'single')
            bottom_border.set(qn('w:sz'), '8')
            borders.append(bottom_border)
            tcPr.append(borders)

        # Adapt cell margins
        tc = table._element
        tblPr = tc.tblPr
        tblCellMar = OxmlElement('w:tblCellMar')
        # set left and right margins to zero
        # and top and bottom margins to 60 dxa
        kwargs = {"left": 0, "right": 0, "top": 60, "bottom": 60}
        for m in ["left", "right", "top", "bottom"]:
            node = OxmlElement("w:{}".format(m))
            node.set(qn('w:w'), str(kwargs.get(m)))
            node.set(qn('w:type'), 'dxa')
            tblCellMar.append(node)
        tblPr.append(tblCellMar)

        # Save the document
        if file_name is not None:
            document.save(file_name)

        return document

    def _output_tex(self, file_name: Optional[str] = None, full_width: bool = False, **kwargs):
        # Make a copy of the DataFrame to avoid modifying the original
        dfs = self.df.copy()

        # First wrap all cells which contain a line break in a makecell command
        dfs = dfs.map(lambda x: f"\\makecell{{{x}}}" if isinstance(x, str) and "\\\\" in x else x)
        row_levels = dfs.index.nlevels
        # when the row index has more than one level, we will store
        # the top level to use later to add clines and row group titles
        # and then remove it
        if row_levels > 1:
            # Store the top level of the row index
            top_row_id = dfs.index.get_level_values(0).to_list()
            # Generate a list of the distinct values
            row_groups = list(dict.fromkeys(top_row_id))
            # Generate a list containing the number of rows for each group
            row_groups_len = [top_row_id.count(group) for group in row_groups]
            # Drop the top level of the row index:
            dfs.index = dfs.index.droplevel(0)

        # Style the table
        styler = dfs.style

        # Generate LaTeX code
        latex_res = styler.to_latex(
            hrules=True,
            multicol_align="c",
            multirow_align="t",
            column_format="l" + "c" * (dfs.shape[1] + dfs.index.nlevels),
        )

        # First split the LaTeX code into lines
        lines = latex_res.splitlines()
        # Find the line number of the \midrule
        line_at = next(i for i, line in enumerate(lines) if "\\midrule" in line)
        # Add space after this \midrule:
        lines.insert(line_at + 1, "\\addlinespace")
        line_at += 1

        # When there are row groups then insert midrules and groupname
        if row_levels > 1 and len(row_groups) > 1:
            # Insert a midrule after each row group
            for i in range(len(row_groups)):
                if self.rgroup_display:
                    # Insert a line with the row group name & same space around it
                    lines.insert(line_at + 1, "\\emph{" + row_groups[i] + "} \\\\")
                    lines.insert(line_at + 2, "\\addlinespace")
                    lines.insert(line_at + 3 + row_groups_len[i], "\\addlinespace")
                    line_at += 3
                if (self.rgroup_sep != "") and (i < len(row_groups) - 1):
                    # For tex output we only either at a line between the row groups or not
                    # And we don't add a line after the last row group
                    line_at += row_groups_len[i] + 1
                    lines.insert(line_at, "\\midrule")
                    lines.insert(line_at + 1, "\\addlinespace")
                    line_at += 1
        else:
            # Add line space before the end of the table
            lines.insert(line_at + dfs.shape[0] + 1, "\\addlinespace")

        # Insert cmidrules (equivalent to column spanners in gt)
        # First find the first line with an occurrence of "multicolumn"
        cmidrule_line_number = None
        for i, line in enumerate(lines):
            if "multicolumn" in line:
                cmidrule_line_number = i + 1
                # Regular expression to find \multicolumn{number}
                pattern = r"\\multicolumn\{(\d+)\}"
                # Find all matches (i.e. values of d) in the LaTeX string & convert to integers
                ncols = [int(match) for match in re.findall(pattern, line)]

                cmidrule_string = ""
                leftcol = 2
                for n in ncols:
                    cmidrule_string += (
                        r"\cmidrule(lr){" + str(leftcol) + "-" + str(leftcol + n - 1) + "} "
                    )
                    leftcol += n
                lines.insert(cmidrule_line_number, cmidrule_string)

        # Put the lines back together
        latex_res = "\n".join(lines)

        # Wrap in threeparttable to allow for table notes
        if self.notes is not None:
            latex_res = (
                "\\begin{threeparttable}\n"
                + latex_res
                + "\n\\footnotesize "
                + self.notes
                + "\n\\end{threeparttable}"
            )
        else:
            latex_res = "\\begin{threeparttable}\n" + latex_res + "\n\\end{threeparttable}"

        # If caption or label specified then wrap in table environment
        if (self.caption is not None) or (self.tab_label is not None):
            latex_res = (
                "\\begin{table}[" + kwargs.get("texlocation", "htbp") + "]\n"
                + "\\centering\n"
                + ("\\caption{" + self.caption + "}\n" if self.caption is not None else "")
                + ("\\label{" + self.tab_label + "}\n" if self.tab_label is not None else "")
                + latex_res
                + "\n\\end{table}"
            )

        # Set cell alignment to top
        latex_res = "\\renewcommand\\cellalign{t}\n" + latex_res

        # Set table width to full page width if full_width is True
        if full_width:
            latex_res = latex_res.replace(
                "\\begin{tabular}{l", "\\begin{tabularx}{\\linewidth}{X"
            )
            latex_res = latex_res.replace(
                "\\end{tabular}", "\\end{tabularx}\n \\vspace{3pt}"
            )

        if file_name is not None:
            with open(file_name, "w") as f:
                f.write(latex_res)  # Write the latex code to a file

        return latex_res

    def _output_gt(self, file_name: Optional[str] = None, full_width: bool = False, **kwargs):
        # Make a copy of the DataFrame to avoid modifying the original
        dfs = self.df.copy()

        # GT does not support MultiIndex columns, so we need to flatten the columns
        if isinstance(dfs.columns, pd.MultiIndex):
            # Store labels of the last level of the column index (to use as column names)
            col_names = dfs.columns.get_level_values(-1)
            nlevels = dfs.columns.nlevels
            
            # Assign column numbers to the lowest index level
            col_numbers = list(map(str, range(len(dfs.columns))))
            # Save the whole column index in order to generate table spanner labels later
            dfcols = dfs.columns.to_list()
            # Flatten the column index just numbering the columns
            dfs.columns = pd.Index(col_numbers)
            # Store the mapping of column numbers to column names
            col_dict = dict(zip(col_numbers, col_names))
            # Modify the last elements in each tuple in dfcols
            dfcols = [(t[:-1] + (col_numbers[i],)) for i, t in enumerate(dfcols)]
        else:
            nlevels = 1

        # store row indes and then reset to have the index as columns to be displayed in the table
        rowindex = dfs.index
        dfs.reset_index(inplace=True)

        # Specify the rowname_col and groupname_col
        if isinstance(rowindex, pd.MultiIndex):
            rowname_col = dfs.columns[1]
            groupname_col = dfs.columns[0]
        else:
            rowname_col = dfs.columns[0]
            groupname_col = None

        # Generate the table with GT
        gt = GT(dfs, auto_align=False)

        # When caption is provided, add it to the table
        if self.caption is not None:
            gt = gt.tab_header(title=self.caption).tab_options(table_border_top_style="hidden")

        # Add column spanners based on multiindex
        if nlevels > 1:
            for i in range(nlevels - 1):
                col_spanners = {}
                # Iterate over columns and group them by the labels in the respective level
                for c in dfcols:
                    key = c[i]
                    if key not in col_spanners:
                        col_spanners[key] = []
                    col_spanners[key].append(c[-1])
                for label, columns in col_spanners.items():
                    gt = gt.tab_spanner(label=label, columns=columns, level=nlevels - 1 - i)
                # Restore column names
                gt = gt.cols_label(**col_dict)

        # Customize the table layout
        gt = (
            gt.tab_source_note(self.notes)
            .tab_stub(rowname_col=rowname_col, groupname_col=groupname_col)
            .tab_options(
            table_border_bottom_style="hidden",
            stub_border_style="hidden",
            column_labels_border_top_style="solid",
            column_labels_border_top_color="black",
            column_labels_border_bottom_style="solid",
            column_labels_border_bottom_color="black",
            column_labels_border_bottom_width="0.5px",
            column_labels_vlines_color="white",
            column_labels_vlines_width="0px",
            table_body_border_top_style="solid",
            table_body_border_top_width="0.5px",
            table_body_border_top_color="black",
            table_body_hlines_style="none",
            table_body_vlines_color="white",
            table_body_vlines_width="0px",
            table_body_border_bottom_color="black",
            row_group_border_top_style="solid",
            row_group_border_top_width="0.5px",
            row_group_border_top_color="black",
            row_group_border_bottom_style="solid",
            row_group_border_bottom_width="0.5px",
            row_group_border_bottom_color="black",
            row_group_border_left_color="white",
            row_group_border_right_color="white",
            data_row_padding="4px",
            column_labels_padding="4px",
            )
            .cols_align(align="center")
        )

        # Full page width
        if full_width:
            gt = gt.tab_options(table_width="100%")

        # Customize row group display
        if "t" not in self.rgroup_sep:
            gt = gt.tab_options(row_group_border_top_style="none")
        if "b" not in self.rgroup_sep:
            gt = gt.tab_options(row_group_border_bottom_style="none")
        if not self.rgroup_display:
            gt = gt.tab_options(row_group_font_size="0px", row_group_padding="0px")

        # Save the html code of the table to a file
        if file_name is not None:
            with open(file_name, "w") as f:
                f.write(gt.as_raw_html())

        return gt
    

    # Sample usage
    if __name__ == "__main__":
        # Create a sample DataFrame
        data = {
            "Group": ["A", "A", "B", "B"],
            "Name": ["John", "Alice", "Bob", "Eve"],
            "Score": [90, 85, 88, 92]
        }
        df = pd.DataFrame(data)
        df.set_index(["Group", "Name"], inplace=True)

        # Create a TabOut object
        tab_out = TabOut(df, notes="Sample notes", caption="Sample Table", tab_label="tab:sample")

        # Generate tables in different formats
        tab_out.make(types=["gt", "tex", "docx", "pptx"], path="sample_table")